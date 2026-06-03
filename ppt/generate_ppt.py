#!/usr/bin/env python3
"""
MediCraft PPT - 数字城市风格（深蓝封面 + 艺术大字 + 科技线条）
封面深色 · 内容白底 · 几何线条 · 大气艺术字
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math

class C:
    # 深色系（封面用）
    NAVY = RGBColor(0x0A, 0x1A, 0x3A)
    NAVY_M = RGBColor(0x0F, 0x25, 0x4A)
    NAVY_L = RGBColor(0x15, 0x30, 0x5A)
    NAVY_LINE = RGBColor(0x1A, 0x40, 0x80)

    # 浅色系（内容页）
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    FOG = RGBColor(0xF5, 0xF7, 0xFA)
    ICE = RGBColor(0xEE, 0xF2, 0xFF)
    SNOW = RGBColor(0xF0, 0xF4, 0xF8)

    # 科技蓝
    BLUE = RGBColor(0x3B, 0x82, 0xF6)
    BLUE_S = RGBColor(0x60, 0xA5, 0xFA)
    BLUE_T = RGBColor(0xDB, 0xEA, 0xFE)
    CYAN = RGBColor(0x06, 0xB6, 0xD4)
    CYAN_S = RGBColor(0x67, 0xE8, 0xF9)
    CYAN_T = RGBColor(0xCC, 0xFB, 0xF1)

    # 渐变紫蓝
    PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
    PURPLE_S = RGBColor(0xA7, 0x8B, 0xFA)
    PURPLE_T = RGBColor(0xED, 0xE9, 0xFE)

    # 点缀
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    GREEN_T = RGBColor(0xD1, 0xFA, 0xE5)
    ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
    ORANGE_T = RGBColor(0xFE, 0xF3, 0xC7)
    PINK = RGBColor(0xEC, 0x48, 0x99)
    PINK_T = RGBColor(0xFC, 0xE7, 0xF3)
    ROSE = RGBColor(0xF4, 0x3F, 0x5E)
    ROSE_T = RGBColor(0xFF, 0xE4, 0xE6)

    # 文字
    H1 = RGBColor(0x0F, 0x17, 0x2A)
    H2 = RGBColor(0x1E, 0x29, 0x3B)
    BODY = RGBColor(0x47, 0x55, 0x69)
    SUB = RGBColor(0x94, 0xA3, 0xB8)
    W = RGBColor(0xFF, 0xFF, 0xFF)
    W_DIM = RGBColor(0xA0, 0xB0, 0xC8)
    BDR = RGBColor(0xE2, 0xE8, 0xF0)

# ===== 工具 =====
def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def R(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()

def RR(s, l, t, w, h, c, b=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    if b:
        sh.line.color.rgb = b
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()

def O(s, l, t, sz, c):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.fill.background()

def G(s, l, t, w, h, c1, c2):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    f = sh.fill
    f.gradient()
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].color.rgb = c2
    sh.line.fill.background()

def Glow(s, l, t, w, h, c1, c2):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    f = sh.fill
    f.gradient()
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[0].color.brightness = 0.8
    f.gradient_stops[1].color.rgb = c2
    f.gradient_stops[1].color.brightness = 0.4
    sh.line.fill.background()

def T(s, l, t, w, h, txt, sz=16, c=C.H1, bold=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = a

def ic(s, x, y, sz, txt, bg_c, tc=C.W, fs=18):
    O(s, x, y, sz, bg_c)
    tb = s.shapes.add_textbox(x, y, sz, sz)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(fs)
    p.font.color.rgb = tc
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(fs * 0.2)

def tech_lines(s, x_base=0.5, y_base=0.5):
    """科技线条网络 - 模拟封面左侧的几何线条"""
    # 主斜线（从左下到右上）
    lines = [
        # (起点x, 起点y, 终点x, 终点y, 颜色, 宽度)
        (0.3, 5.5, 4.0, 0.5, C.NAVY_LINE, 0.02),
        (0.5, 6.0, 5.0, 1.0, C.NAVY_LINE, 0.015),
        (1.0, 6.5, 6.0, 1.5, C.NAVY_LINE, 0.012),
        (0.8, 5.0, 3.5, 1.5, C.NAVY_LINE, 0.018),
        (1.5, 6.0, 5.5, 2.0, C.NAVY_LINE, 0.01),
        # 交叉线
        (2.0, 0.3, 5.0, 4.5, C.NAVY_LINE, 0.012),
        (3.0, 0.5, 6.0, 5.0, C.NAVY_LINE, 0.01),
        (1.0, 1.0, 4.5, 6.0, C.NAVY_LINE, 0.008),
    ]

    for (x1, y1, x2, y2, color, w) in lines:
        # 用小圆点模拟线条
        steps = 30
        for j in range(steps + 1):
            t = j / steps
            px = x1 + (x2 - x1) * t
            py = y1 + (y2 - y1) * t
            O(s, Inches(px), Inches(py), Inches(w), color)

    # 节点（交叉点）
    nodes = [
        (2.5, 2.5, 0.12), (3.5, 3.5, 0.1), (1.5, 3.0, 0.08),
        (4.0, 2.0, 0.09), (3.0, 1.5, 0.07), (2.0, 4.0, 0.1),
        (4.5, 3.0, 0.08), (1.0, 4.5, 0.06), (3.0, 5.0, 0.07),
    ]
    for (x, y, sz) in nodes:
        O(s, Inches(x), Inches(y), Inches(sz), C.NAVY_LINE)
        O(s, Inches(x + sz * 0.25), Inches(y + sz * 0.25), Inches(sz * 0.4), C.BLUE)

def tech_lines_light(s):
    """浅色版科技线条（内容页用）"""
    lines = [
        (0.3, 5.5, 4.0, 0.5, C.BLUE_T),
        (0.5, 6.0, 5.0, 1.0, C.PURPLE_T),
        (2.0, 0.3, 5.0, 4.5, C.CYAN_T),
        (3.0, 0.5, 6.0, 5.0, C.BLUE_T),
    ]
    for (x1, y1, x2, y2, color) in lines:
        steps = 25
        for j in range(steps + 1):
            t = j / steps
            px = x1 + (x2 - x1) * t
            py = y1 + (y2 - y1) * t
            O(s, Inches(px), Inches(py), Inches(0.012), color)

    nodes = [
        (2.5, 2.5, 0.08), (3.5, 3.5, 0.06), (1.5, 3.0, 0.05),
        (4.0, 2.0, 0.06), (3.0, 1.5, 0.05),
    ]
    for (x, y, sz) in nodes:
        O(s, Inches(x), Inches(y), Inches(sz), C.BLUE_T)

def glass(s, x, y, w, h):
    RR(s, x, y, w, h, C.WHITE, C.BDR)
    R(s, x + Inches(0.1), y + Inches(0.04), Inches(0.5), Pt(2), C.BLUE_T)

def top_bar(s):
    G(s, Inches(0), Inches(0), Inches(13.333), Pt(4), C.BLUE, C.CYAN)

def pn(s, n, t):
    T(s, Inches(12.1), Inches(7.05), Inches(1), Inches(0.3), f'{n}/{t}', sz=9, c=C.SUB, a=PP_ALIGN.RIGHT)

def sec(s, main, sub="", y=Inches(0.35)):
    T(s, Inches(0.9), y, Inches(10), Inches(0.6), main, sz=32, c=C.H1, bold=True)
    if sub:
        T(s, Inches(0.9), y + Inches(0.55), Inches(10), Inches(0.3), sub, sz=13, c=C.SUB)
    G(s, Inches(0.9), y + Inches(0.95), Inches(1.6), Pt(3), C.BLUE, C.CYAN)

def create():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    N = 16

    # ────────────────────────────────────
    # 1 封面（深色 + 艺术大字 + 科技线条）
    # ────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.NAVY)

    # 科技线条网络（左侧大面积）
    tech_lines(s)

    # 右侧渐变光晕
    Glow(s, Inches(7), Inches(0.5), Inches(7), Inches(6), C.NAVY_M, C.NAVY)

    # 左侧竖线装饰
    R(s, Inches(0.8), Inches(0.5), Pt(3), Inches(6.5), C.BLUE)

    # 超大艺术字标题
    T(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
      'AI 智能学习辅助系统', sz=60, c=C.W, bold=True)

    # 横线
    G(s, Inches(1.5), Inches(3.2), Inches(5), Pt(3), C.BLUE, C.CYAN)

    # 英文副标题
    T(s, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
      'AI INTELLIGENT LEARNING ASSISTANCE SYSTEM', sz=20, c=C.W_DIM)

    # 项目描述
    T(s, Inches(1.5), Inches(4.3), Inches(8), Inches(0.5),
      '基于多智能体协同的个性化学习路径规划与效果评估平台', sz=15, c=C.W_DIM)

    # 底部标签
    tags = ['Spring Boot 3', 'Vue 3', 'Spring AI', '多智能体', '个性化学习']
    x = Inches(1.5)
    for tag in tags:
        RR(s, x, Inches(5.2), Inches(1.4), Inches(0.33), C.NAVY_L, C.NAVY_LINE)
        T(s, x, Inches(5.22), Inches(1.4), Inches(0.3), tag, sz=10, c=C.W_DIM, a=PP_ALIGN.CENTER)
        x += Inches(1.55)

    # 底部信息
    T(s, Inches(1.5), Inches(6.2), Inches(4), Inches(0.25),
      'MediCraft Development Team', sz=12, c=C.SUB)
    T(s, Inches(1.5), Inches(6.5), Inches(4), Inches(0.25),
      '2026 年度大学生创新创业项目', sz=11, c=C.SUB)

    # ────────────────────────────────────
    # 2 目录（白底 + 浅色科技线条）
    # ────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.WHITE)
    tech_lines_light(s)
    top_bar(s)
    sec(s, '目录', 'CONTENTS')

    toc = [
        ('01', '项目背景', '传统教育痛点与AI赋能机遇', C.BLUE),
        ('02', '系统概述', '全链路AI辅助学习平台', C.CYAN),
        ('03', '核心功能', '五大智能模块详解', C.PURPLE),
        ('04', '技术架构', '前后端分离+多智能体协同', C.GREEN),
        ('05', '创新亮点', '架构创新与工程实践', C.ORANGE),
        ('06', '项目总结', '成果展示与未来规划', C.PINK),
    ]
    for i, (num, t, d, c) in enumerate(toc):
        y = Inches(1.5) + i * Inches(0.9)
        ic(s, Inches(1.3), y, Inches(0.52), num, c, fs=16)
        T(s, Inches(2.1), y + Inches(0.02), Inches(3.5), Inches(0.3), t, sz=20, c=C.H2, bold=True)
        T(s, Inches(2.1), y + Inches(0.35), Inches(4), Inches(0.2), d, sz=11, c=C.SUB)
        if i < 5:
            G(s, Inches(2.1), y + Inches(0.68), Inches(9.5), Pt(1), C.BLUE_T, C.CYAN_T)
    pn(s, 2, N)

    # ────────────────────────────────────
    # 3-16 内容页（白底 + 浅色科技线条装饰）
    # ────────────────────────────────────
    content_pages = [
        # (背景色, 标题, 副标题)
        (C.WHITE, '项目背景', 'WHY WE BUILD THIS'),
        (C.FOG, '系统概述', 'MediCraft - AI SMART LEARNING PLATFORM'),
        (C.WHITE, '模块一：用户与学习画像', 'AI-POWERED STUDENT PROFILING'),
        (C.WHITE, '模块二：AI 多智能体资源生成', 'MULTI-AGENT RESOURCE GENERATION'),
        (C.WHITE, '资源生成流程', 'ASYNCHRONOUS GENERATION PIPELINE'),
        (C.FOG, '模块三：个性化学习路径规划', 'DYNAMIC LEARNING PATH PLANNING'),
        (C.WHITE, '模块四：智能辅导答疑', 'AI-POWERED INTELLIGENT TUTORING'),
        (C.WHITE, '模块五：学习效果评估', 'MULTI-DIMENSIONAL LEARNING ASSESSMENT'),
        (C.WHITE, '技术架构', 'TECHNICAL ARCHITECTURE'),
        (C.FOG, '技术栈详情', 'TECHNOLOGY STACK'),
        (C.WHITE, '数据库设计', 'DATABASE DESIGN - 14 TABLES'),
        (C.WHITE, '创新亮点', 'INNOVATION HIGHLIGHTS'),
        (C.WHITE, '项目成果与未来规划', 'ACHIEVEMENTS & ROADMAP'),
    ]

    for idx, (bg_c, title_t, sub_t) in enumerate(content_pages):
        pg = idx + 3
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg(s, bg_c)
        tech_lines_light(s)
        top_bar(s)
        sec(s, title_t, sub_t)

        # 根据页码渲染不同内容
        if pg == 3:  # 项目背景
            for i, (hd, items, accent) in enumerate([
                ('传统教育困境', ['学习资源千篇一律，无法满足个性化需求', '学生缺乏科学的学习路径规划指导', '学习效果评估依赖考试，维度单一', '师生比失衡，个性化辅导难以实现'], C.ROSE),
                ('AI 赋能机遇', ['大语言模型可实现智能内容生成', '多智能体协同处理复杂教育任务', '数据驱动的精准学情分析成为可能', '个性化学习路径自动规划与调整'], C.GREEN),
            ]):
                x = Inches(0.8) + i * Inches(6.2)
                glass(s, x, Inches(1.5), Inches(5.8), Inches(5))
                R(s, x + Inches(0.18), Inches(1.58), Inches(0.8), Pt(3), accent)
                T(s, x + Inches(0.35), Inches(1.75), Inches(4), Inches(0.35), hd, sz=20, c=C.H2, bold=True)
                for j, item in enumerate(items):
                    y = Inches(2.4) + j * Inches(0.85)
                    O(s, x + Inches(0.4), y + Inches(0.06), Inches(0.13), accent)
                    T(s, x + Inches(0.72), y, Inches(4.5), Inches(0.6), item, sz=13, c=C.BODY)

        elif pg == 4:  # 系统概述
            T(s, Inches(0.9), Inches(1.45), Inches(11), Inches(0.5),
              'MediCraft 是一个基于多智能体协同的个性化学习平台，为学生提供从画像构建、资源生成、路径规划到效果评估的全链路 AI 辅助学习体验。',
              sz=14, c=C.BODY)
            mods = [('学习画像', '6维画像构建', C.BLUE), ('资源生成', '6大AI智能体', C.CYAN),
                    ('路径规划', '动态路径调整', C.PURPLE), ('智能辅导', '多模态答疑', C.GREEN),
                    ('效果评估', '多维度评估', C.ORANGE)]
            Glow(s, Inches(5.3), Inches(2.7), Inches(2.3), Inches(2.3), C.BLUE_T, C.CYAN_T)
            ic(s, Inches(5.9), Inches(3.2), Inches(1.2), 'AI', C.BLUE, fs=30)
            pos = [(Inches(1.3), Inches(2.6)), (Inches(3.6), Inches(4.8)),
                   (Inches(5.9), Inches(5.2)), (Inches(8.2), Inches(4.8)),
                   (Inches(10.5), Inches(2.6))]
            for i, ((t, d, c), (x, y)) in enumerate(zip(mods, pos)):
                glass(s, x, y, Inches(2.1), Inches(1.4))
                R(s, x + Inches(0.12), y + Inches(0.06), Inches(0.4), Pt(3), c)
                ic(s, x + Inches(0.65), y + Inches(0.15), Inches(0.42), str(i + 1), c, fs=14)
                T(s, x + Inches(0.1), y + Inches(0.65), Inches(1.9), Inches(0.25), t, sz=14, c=C.H2, bold=True, a=PP_ALIGN.CENTER)
                T(s, x + Inches(0.1), y + Inches(0.95), Inches(1.9), Inches(0.2), d, sz=10, c=C.SUB, a=PP_ALIGN.CENTER)

        elif pg == 5:  # 用户画像
            glass(s, Inches(0.8), Inches(1.5), Inches(6), Inches(5.3))
            T(s, Inches(1.2), Inches(1.7), Inches(5), Inches(0.35), '对话式画像构建', sz=18, c=C.BLUE, bold=True)
            T(s, Inches(1.2), Inches(2.1), Inches(5), Inches(0.25), '通过多轮 AI 对话，自然收集学生学习特征', sz=11, c=C.SUB)
            for i, (role, msg, bg_c) in enumerate([('AI', '你好！请告诉我你的专业和年级？', C.BLUE_T), ('用户', '我是计算机科学专业，大三学生', C.GREEN_T), ('AI', '你的学习目标是什么？希望提升哪些方面？', C.BLUE_T), ('用户', '想系统学习Spring Boot，准备春招', C.GREEN_T)]):
                y = Inches(2.6) + i * Inches(0.68)
                x = Inches(1.2) if role == 'AI' else Inches(1.8)
                RR(s, x, y, Inches(4.6), Inches(0.5), C.FOG, C.BDR)
                T(s, x + Inches(0.1), y + Inches(0.1), Inches(4.3), Inches(0.3), f'{role}: {msg}', sz=11, c=C.BODY)
            glass(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.3))
            T(s, Inches(7.6), Inches(1.7), Inches(4), Inches(0.35), '六维学习画像', sz=18, c=C.BLUE, bold=True)
            for i, (d, desc, c) in enumerate([('知识基础', '已有知识水平评估', C.BLUE), ('认知风格', '学习偏好分析', C.CYAN), ('学习目标', '阶段性目标设定', C.PURPLE), ('易错点', '薄弱知识点识别', C.ORANGE), ('学习节奏', '学习强度适配', C.GREEN), ('资源偏好', '内容形式偏好', C.PINK)]):
                r, co = i // 2, i % 2
                x = Inches(7.5) + co * Inches(2.4)
                y = Inches(2.3) + r * Inches(1.35)
                RR(s, x, y, Inches(2.15), Inches(1.1), C.FOG, C.BDR)
                R(s, x + Inches(0.1), y + Inches(0.06), Inches(0.4), Pt(3), c)
                T(s, x + Inches(0.12), y + Inches(0.18), Inches(1.9), Inches(0.25), d, sz=13, c=C.H2, bold=True)
                T(s, x + Inches(0.12), y + Inches(0.5), Inches(1.9), Inches(0.2), desc, sz=10, c=C.SUB)

        elif pg == 6:  # 多智能体
            for i, (nm, role, desc, c) in enumerate([('需求解析', 'demand', '学习目标拆解 / 知识清单 / 推荐路径', C.BLUE), ('文档生成', 'document', '结构化文档 / 学习笔记 / 知识详解', C.CYAN), ('思维导图', 'mind', '知识体系 / 关联可视化 / 复习提纲', C.PURPLE), ('题库生成', 'question', '多难度练习 / 选择填空 / 解析评分', C.ORANGE), ('实操案例', 'case', '代码实战 / 项目拆解 / 调试排错', C.GREEN), ('多模态', 'multimodal', '视频脚本 / 音频讲解 / 动画文案', C.PINK)]):
                r, co = i // 3, i % 3
                x = Inches(0.8) + co * Inches(4.1)
                y = Inches(1.5) + r * Inches(2.8)
                glass(s, x, y, Inches(3.8), Inches(2.5))
                ic(s, x + Inches(0.18), y + Inches(0.2), Inches(0.46), str(i + 1), c, fs=16)
                T(s, x + Inches(0.78), y + Inches(0.25), Inches(2.5), Inches(0.25), f'{nm}智能体', sz=15, c=C.H2, bold=True)
                T(s, x + Inches(0.78), y + Inches(0.5), Inches(2.5), Inches(0.2), f'Role: {role}', sz=9, c=c)
                T(s, x + Inches(0.3), y + Inches(0.95), Inches(3.2), Inches(1.2), desc, sz=12, c=C.BODY)

        elif pg == 7:  # 生成流程
            for i, (num, t, d, c) in enumerate([('01', '选择智能体', '从6大智能体中\n选择目标角色', C.BLUE), ('02', '输入主题', '填写学习主题\n难度与参数', C.CYAN), ('03', '异步生成', '后端异步调用AI\n实时进度轮询', C.PURPLE), ('04', '质量校验', 'AI内容格式校验\n自动修复兼容', C.GREEN), ('05', '资源入库', '结果持久化\n关联学习路径', C.ORANGE)]):
                x = Inches(0.5) + i * Inches(2.5)
                y = Inches(1.6)
                glass(s, x, y, Inches(2.2), Inches(2.5))
                ic(s, x + Inches(0.7), y + Inches(0.18), Inches(0.52), num, c, fs=16)
                T(s, x + Inches(0.08), y + Inches(0.85), Inches(2.05), Inches(0.25), t, sz=15, c=C.H2, bold=True, a=PP_ALIGN.CENTER)
                for j, line in enumerate(d.split('\n')):
                    T(s, x + Inches(0.08), y + Inches(1.2) + j * Inches(0.26), Inches(2.05), Inches(0.22), line, sz=10.5, c=C.BODY, a=PP_ALIGN.CENTER)
                if i < 4:
                    T(s, x + Inches(2.12), y + Inches(0.85), Inches(0.5), Inches(0.3), '>', sz=20, c=c, bold=True, a=PP_ALIGN.CENTER)
            glass(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3))
            T(s, Inches(1.2), Inches(4.65), Inches(3), Inches(0.3), '核心技术要点', sz=16, c=C.BLUE, bold=True)
            for i, p in enumerate(['策略模式: AgentStrategyFactory 自动发现 + 按 role 分发，新增智能体零改动', '异步解耦: ResourceAsyncExecutor 独立 Bean，解决 Spring @Async 自调用失效', 'Prompt 工程化: 每个智能体独立 Prompt 模板，存储在数据库支持动态配置', '进度轮询: 前端 2s 间隔轮询，实时展示 0-100% 生成进度']):
                T(s, Inches(1.2), Inches(5.05) + i * Inches(0.28), Inches(10.5), Inches(0.25), p, sz=11, c=C.BODY)

        elif pg == 8:  # 学习路径
            glass(s, Inches(0.8), Inches(1.5), Inches(6), Inches(5.3))
            T(s, Inches(1.2), Inches(1.7), Inches(5), Inches(0.35), 'AI 智能路径生成', sz=18, c=C.GREEN, bold=True)
            for i, (nm, desc, c) in enumerate([('基础语法', '变量、数据类型、控制流', C.GREEN), ('框架入门', 'Spring Boot 核心概念', C.BLUE), ('数据库', 'MySQL + MyBatis-Plus', C.PURPLE), ('项目实战', '完整项目开发与部署', C.CYAN), ('知识测验', '综合测验与评估', C.PINK)]):
                y = Inches(2.4) + i * Inches(0.8)
                O(s, Inches(1.35), y + Inches(0.04), Inches(0.2), c)
                if i < 4: R(s, Inches(1.42), y + Inches(0.24), Pt(1.2), Inches(0.58), C.BDR)
                T(s, Inches(1.75), y, Inches(2), Inches(0.22), nm, sz=12, c=C.H2, bold=True)
                T(s, Inches(1.75), y + Inches(0.25), Inches(3.5), Inches(0.2), desc, sz=10, c=C.SUB)
            glass(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.3))
            T(s, Inches(7.6), Inches(1.7), Inches(4), Inches(0.35), '5 种动态调整方式', sz=18, c=C.ORANGE, bold=True)
            for i, (nm, desc, c) in enumerate([('延长周期', '学习节奏太快，需要更多时间', C.BLUE), ('压缩周期', '学习进度超前，加快节奏', C.CYAN), ('调整顺序', '拖拽调整学习内容先后顺序', C.PURPLE), ('新增内容', '补充额外知识点到路径中', C.GREEN), ('简化内容', '移除已掌握的学习步骤', C.PINK)]):
                y = Inches(2.4) + i * Inches(0.82)
                RR(s, Inches(7.5), y, Inches(4.6), Inches(0.65), C.FOG, C.BDR)
                R(s, Inches(7.58), y + Inches(0.06), Inches(0.35), Pt(3), c)
                T(s, Inches(7.68), y + Inches(0.08), Inches(1.2), Inches(0.22), nm, sz=12, c=c, bold=True)
                T(s, Inches(9.1), y + Inches(0.08), Inches(2.8), Inches(0.5), desc, sz=10.5, c=C.BODY)

        elif pg == 9:  # 智能辅导
            glass(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
            T(s, Inches(1.3), Inches(1.7), Inches(5), Inches(0.35), '多模态 AI 答疑', sz=18, c=C.ORANGE, bold=True)
            for i, (role, msg, bg_c) in enumerate([('user', '什么是 Spring Boot 的自动配置原理？', C.GREEN_T), ('ai', 'Spring Boot 自动配置基于 @EnableAutoConfiguration...', C.BLUE_T), ('user', '能举个具体的例子吗？', C.GREEN_T), ('ai', '当然！以 DataSourceAutoConfiguration 为例...', C.BLUE_T)]):
                y = Inches(2.4) + i * Inches(0.68)
                x = Inches(1.4) if role == 'ai' else Inches(2.0)
                RR(s, x, y, Inches(5), Inches(0.52), bg_c, C.BDR)
                T(s, x + Inches(0.1), y + Inches(0.1), Inches(4.7), Inches(0.3), f'{"AI" if role == "ai" else "我"}: {msg}', sz=11, c=C.BODY)
            for i, (nm, desc) in enumerate([('文字提问', '输入问题获取详细解答'), ('图片识别', '支持截图/拍照识别题目'), ('逐字显示', '打字机效果流式输出'), ('历史记录', '保存所有答疑历史'), ('多轮对话', '支持追问和深入探讨')]):
                y = Inches(2.4) + i * Inches(0.68)
                O(s, Inches(7.8), y + Inches(0.04), Inches(0.14), C.ORANGE)
                T(s, Inches(8.08), y, Inches(2), Inches(0.22), nm, sz=12, c=C.H2, bold=True)
                T(s, Inches(8.08), y + Inches(0.26), Inches(3), Inches(0.2), desc, sz=10, c=C.SUB)

        elif pg == 10:  # 效果评估
            for i, (hd, items, c) in enumerate([('AI 评估报告', ['根据学习行为数据生成', '多维度评估分析', '知识点掌握度 JSON', '薄弱点精准识别', '个性化提升建议'], C.BLUE), ('数据统计', ['总学习时长统计', '答题正确率趋势', '每日学习数据', '资源类型分布', '各学科进度对比'], C.CYAN), ('能力雷达', ['知识掌握度评估', '技能应用能力', '学习进度追踪', '练习准确率分析', '学习一致性评分'], C.PURPLE)]):
                x = Inches(0.8) + i * Inches(4.1)
                glass(s, x, Inches(1.5), Inches(3.8), Inches(5.3))
                T(s, x + Inches(0.3), Inches(1.75), Inches(3), Inches(0.35), hd, sz=18, c=c, bold=True)
                for j, item in enumerate(items):
                    y = Inches(2.4) + j * Inches(0.78)
                    O(s, x + Inches(0.35), y + Inches(0.05), Inches(0.13), c)
                    T(s, x + Inches(0.62), y, Inches(2.8), Inches(0.5), item, sz=12, c=C.BODY)

        elif pg == 11:  # 技术架构
            for i, (nm, tech, mods, c) in enumerate([('前端层', 'Vue 3 + Vite + Element Plus + Pinia + Vue Router + Axios', ['登录注册', '画像构建', '资源生成', '路径规划', '智能辅导', '效果评估'], C.BLUE), ('后端层', 'Spring Boot 3 + Spring AI + MyBatis-Plus + JWT + Redis', ['用户模块', '画像模块', '资源模块', '路径模块', '辅导模块', '评估模块'], C.CYAN), ('数据层', 'MySQL 8.0 + Redis 6.0 + 14张业务表', ['用户表', '画像表', '资源表', '路径表', '行为表', '评估表'], C.PURPLE)]):
                y = Inches(1.5) + i * Inches(1.9)
                glass(s, Inches(0.8), y, Inches(11.5), Inches(1.65))
                R(s, Inches(0.88), y + Inches(0.06), Inches(0.6), Pt(3), c)
                T(s, Inches(1.0), y + Inches(0.18), Inches(1.2), Inches(0.3), nm, sz=16, c=c, bold=True)
                T(s, Inches(2.3), y + Inches(0.2), Inches(9), Inches(0.25), tech, sz=11, c=C.SUB)
                for j, mod in enumerate(mods):
                    mx = Inches(1.0) + j * Inches(1.8)
                    RR(s, mx, y + Inches(0.6), Inches(1.6), Inches(0.7), C.FOG, C.BDR)
                    T(s, mx, y + Inches(0.72), Inches(1.6), Inches(0.35), mod, sz=11, c=C.BODY, a=PP_ALIGN.CENTER)

        elif pg == 12:  # 技术栈
            glass(s, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.4))
            T(s, Inches(1.2), Inches(1.55), Inches(3), Inches(0.35), '后端技术栈', sz=18, c=C.BLUE, bold=True)
            for i, (tech, desc, c) in enumerate([('Spring Boot 3.3.5', '应用框架', C.BLUE), ('Java 17', '编程语言', C.CYAN), ('Spring AI 1.0.0-M6', 'AI 框架', C.PURPLE), ('MyBatis-Plus 3.5.11', 'ORM 框架', C.GREEN), ('MySQL 8.0+', '数据库', C.ORANGE), ('Redis 6.0+', '缓存', C.PINK), ('JWT', '用户认证', C.ROSE)]):
                y = Inches(2.1) + i * Inches(0.58)
                O(s, Inches(1.3), y + Inches(0.06), Inches(0.12), c)
                T(s, Inches(1.58), y, Inches(2.3), Inches(0.22), tech, sz=12, c=C.H2, bold=True)
                T(s, Inches(4.1), y, Inches(2), Inches(0.22), desc, sz=10, c=C.SUB)
            glass(s, Inches(7), Inches(1.4), Inches(5.5), Inches(5.4))
            T(s, Inches(7.4), Inches(1.55), Inches(3), Inches(0.35), '前端技术栈', sz=18, c=C.BLUE, bold=True)
            for i, (tech, desc, c) in enumerate([('Vue 3', '前端框架', C.BLUE), ('Vite 5', '构建工具', C.CYAN), ('Element Plus', 'UI 组件库', C.PURPLE), ('Pinia', '状态管理', C.GREEN), ('Vue Router', '前端路由', C.ORANGE), ('Axios', 'HTTP 请求', C.PINK), ('Markdown 渲染', '富文本展示', C.ROSE)]):
                y = Inches(2.1) + i * Inches(0.58)
                O(s, Inches(7.5), y + Inches(0.06), Inches(0.12), c)
                T(s, Inches(7.78), y, Inches(2.3), Inches(0.22), tech, sz=12, c=C.H2, bold=True)
                T(s, Inches(10.3), y, Inches(2), Inches(0.22), desc, sz=10, c=C.SUB)

        elif pg == 13:  # 数据库
            for i, (nm, desc, c) in enumerate([('sys_user', '用户表', C.BLUE), ('student_profile', '学习画像', C.CYAN), ('chat_context', '对话上下文', C.PURPLE), ('ai_agent', 'AI智能体', C.GREEN), ('learning_resource', '学习资源', C.BLUE), ('learning_path', '学习路径', C.CYAN), ('learning_path_step', '路径步骤', C.PURPLE), ('step_resource', '步骤资源', C.GREEN), ('learning_behavior', '学习行为', C.ORANGE), ('question_answer', '答题记录', C.PINK), ('video_progress', '视频进度', C.ROSE), ('learning_evaluate', '学习评估', C.ORANGE), ('smart_tutor', '智能辅导', C.PINK), ('generate_task', '生成任务', C.ROSE)]):
                r, co = i // 4, i % 4
                x = Inches(0.6) + co * Inches(3.1)
                y = Inches(1.5) + r * Inches(1.8)
                glass(s, x, y, Inches(2.85), Inches(1.5))
                T(s, x + Inches(0.18), y + Inches(0.28), Inches(2.5), Inches(0.25), nm, sz=11.5, c=c, bold=True)
                T(s, x + Inches(0.18), y + Inches(0.6), Inches(2.5), Inches(0.25), desc, sz=11, c=C.SUB)

        elif pg == 14:  # 创新亮点
            for i, (t, d, c) in enumerate([('多智能体策略模式', '6个智能体各自独立策略类\n新增智能体只需实现接口\n符合开闭原则 OCP', C.BLUE), ('对话式画像构建', 'AI多轮对话收集信息\n自动抽取结构化字段\n替代传统表单填写', C.CYAN), ('异步任务解耦', '独立Bean解决@Async失效\n前端轮询实时进度\n零长连接依赖', C.PURPLE), ('Prompt 工程化', '各智能体独立Prompt模板\n数据库存储动态配置\n结构化输出确保质量', C.GREEN), ('动态路径调整', '5种调整方式\nAI自动重新规划\n实时更新倒计时', C.ORANGE), ('全链路评估', '行为+答题+AI分析\n多维度能力雷达\n自动同步画像更新', C.PINK)]):
                r, co = i // 3, i % 3
                x = Inches(0.6) + co * Inches(4.15)
                y = Inches(1.5) + r * Inches(2.8)
                glass(s, x, y, Inches(3.85), Inches(2.5))
                ic(s, x + Inches(0.18), y + Inches(0.18), Inches(0.42), str(i + 1), c, fs=14)
                T(s, x + Inches(0.72), y + Inches(0.22), Inches(2.8), Inches(0.3), t, sz=15, c=C.H2, bold=True)
                for j, line in enumerate(d.split('\n')):
                    T(s, x + Inches(0.28), y + Inches(0.75) + j * Inches(0.33), Inches(3.3), Inches(0.28), line, sz=11, c=C.BODY)

        elif pg == 15:  # 成果与规划
            glass(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3))
            T(s, Inches(1.2), Inches(1.65), Inches(3), Inches(0.35), '项目成果', sz=18, c=C.GREEN, bold=True)
            for i, (num, label, c) in enumerate([('5', '核心功能模块', C.BLUE), ('6', 'AI 智能体角色', C.CYAN), ('14', '数据库表设计', C.PURPLE)]):
                x = Inches(1.2) + i * Inches(1.8)
                T(s, x, Inches(2.3), Inches(1.5), Inches(0.45), num, sz=32, c=c, bold=True, a=PP_ALIGN.CENTER)
                T(s, x, Inches(2.72), Inches(1.5), Inches(0.2), label, sz=10, c=C.SUB, a=PP_ALIGN.CENTER)
            for i, (num, label, c) in enumerate([('22', '前端页面组件', C.GREEN), ('19', 'REST API 端点', C.ORANGE)]):
                x = Inches(1.2) + i * Inches(1.8)
                T(s, x, Inches(3.5), Inches(1.5), Inches(0.45), num, sz=32, c=c, bold=True, a=PP_ALIGN.CENTER)
                T(s, x, Inches(3.92), Inches(1.5), Inches(0.2), label, sz=10, c=C.SUB, a=PP_ALIGN.CENTER)
            glass(s, Inches(7), Inches(1.5), Inches(5.5), Inches(5.3))
            T(s, Inches(7.4), Inches(1.65), Inches(3), Inches(0.35), '未来规划', sz=18, c=C.BLUE, bold=True)
            for i, (t, d, c) in enumerate([('SSE 流式输出', '接入 Server-Sent Events 实现真正流式响应', C.BLUE), ('Redis 深度缓存', '画像、资源、统计数据全面缓存优化', C.CYAN), ('WebSocket 通知', '实时推送通知（新资源、路径建议）', C.PURPLE), ('加权评估算法', '多维度加权评分，更精准的效果评估', C.GREEN), ('单元测试覆盖', '补充 Service/Controller 层测试用例', C.ORANGE)]):
                y = Inches(2.2) + i * Inches(0.85)
                O(s, Inches(7.5), y + Inches(0.06), Inches(0.14), c)
                T(s, Inches(7.78), y, Inches(4), Inches(0.22), t, sz=12, c=C.H2, bold=True)
                T(s, Inches(7.78), y + Inches(0.28), Inches(4.3), Inches(0.35), d, sz=10, c=C.SUB)

        pn(s, pg, N)

    # ────────────────────────────────────
    # 16 致谢（深色封面风格）
    # ────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.NAVY)
    tech_lines(s)
    Glow(s, Inches(3), Inches(1), Inches(7), Inches(2), C.NAVY_M, C.NAVY)

    R(s, Inches(0.8), Inches(0.5), Pt(3), Inches(6.5), C.BLUE)

    T(s, Inches(0), Inches(1.8), Inches(13.333), Inches(1.2),
      'THANK YOU', sz=64, c=C.W, bold=True, a=PP_ALIGN.CENTER)
    G(s, Inches(5.5), Inches(3.1), Inches(2.5), Pt(3), C.BLUE, C.CYAN)
    T(s, Inches(0), Inches(3.5), Inches(13.333), Inches(0.5),
      'MediCraft - AI 智能学习辅助系统', sz=22, c=C.BLUE_S, a=PP_ALIGN.CENTER)
    T(s, Inches(0), Inches(4.2), Inches(13.333), Inches(0.4),
      '基于多智能体协同的个性化学习路径规划与效果评估平台', sz=14, c=C.W_DIM, a=PP_ALIGN.CENTER)

    tags2 = ['Spring Boot 3', 'Vue 3', 'Spring AI', '多智能体', '个性化学习']
    tw = len(tags2) * 1.45 + (len(tags2) - 1) * 0.12
    x = (13.333 - tw) / 2
    for tag in tags2:
        RR(s, Inches(x), Inches(4.9), Inches(1.45), Inches(0.33), C.NAVY_L, C.NAVY_LINE)
        T(s, Inches(x), Inches(4.92), Inches(1.45), Inches(0.3), tag, sz=10, c=C.W_DIM, a=PP_ALIGN.CENTER)
        x += 1.57
    T(s, Inches(0), Inches(5.6), Inches(13.333), Inches(0.3),
      '2026 年度大学生创新创业项目', sz=12, c=C.SUB, a=PP_ALIGN.CENTER)
    pn(s, 16, N)

    out = r'C:\Users\57780\Desktop\code2\CAL\ppt\MediCraft_v7.pptx'
    prs.save(out)
    print(f'Done: {out}')

if __name__ == '__main__':
    create()
